# -*- coding: utf-8 -*-
"""
Created on Thu May 23 10:27:44 2024

@author: piretro999 Paolo Forte
"""

import os
import pandas as pd
from PyPDF2 import PdfReader
import fitz  # PyMuPDF, utilizzato per file PDF
from pptx import Presentation  # per la lettura dei file PowerPoint
from moviepy.editor import VideoFileClip
import speech_recognition as sr
from docx import Document
import xml.etree.ElementTree as ET
from pydub import AudioSegment
import csv
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
import zipfile
import tempfile
import shutil

temp_directory = r'L:\\temp'  # Temporary directory for file processing

def remove_headers_footers(text):
    """Removes headers and footers from extracted text."""
    lines = text.split('\n')
    if len(lines) > 3:
        return '\n'.join(lines[1:-1])
    return text

def handle_text_file(file_path):
    """Manages text files, removing headers and page feet."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            text = file.read()
        return remove_headers_footers(text), file_path
    except Exception as e:
        return f"Failed to read or process text file: {str(e)}", None

def handle_pdf_file(file_path):
    """Extracts text from PDF files using PyMuPDF."""
    try:
        doc = fitz.open(file_path)
        text = [page.get_text("text") for page in doc]
        doc.close()
        return remove_headers_footers('\n'.join(text)), file_path
    except Exception as e:
        return f"Failed to process PDF file: {str(e)}", None

def handle_word_file(file_path):
    """Extracts text from Word (.docx, .doc) files using python-docx."""
    try:
        doc = Document(file_path)
        text = '\n'.join([para.text for para in doc.paragraphs])
        return remove_headers_footers(text), file_path
    except Exception as e:
        return f"Failed to process Word file: {str(e)}", None

def handle_ppt_file(file_path):
    """Extracts text from PowerPoint presentations (.pptx, .ppt)."""
    try:
        ppt = Presentation(file_path)
        text = [shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")]
        return remove_headers_footers('\n'.join(text)), file_path
    except Exception as e:
        return f"Failed to process PowerPoint file: {str(e)}", None

def handle_excel_file(file_path):
    """Converts Excel spreadsheets to CSV format."""
    try:
        df = pd.read_excel(file_path)
        return df.to_csv(index=False), file_path
    except Exception as e:
        return f"Failed to process Excel file: {str(e)}", None

def handle_csv_file(file_path):
    """Reads CSV files and transforms them back into standardized CSV strings."""
    try:
        with open(file_path, mode='r', encoding='utf-8', newline='') as f:
            reader = csv.reader(f)
            data = list(reader)
        return '\n'.join([','.join(row) for row in data]), file_path
    except Exception as e:
        return f"Failed to process CSV file: {str(e)}", None

def handle_epub_file(file_path):
    """Reads and cleans HTML text from EPUB files."""
    try:
        book = epub.read_epub(file_path)
        text = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text.append(soup.get_text())
        return remove_headers_footers('\n'.join(text)), file_path
    except Exception as e:
        return f"Failed to process EPUB file: {str(e)}", None

def handle_xml_gan_file(file_path):
    """Extracts text from XML and GAN files using ElementTree."""
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        texts = [elem.text for elem in root.iter() if elem.text is not None]
        return '\n'.join(texts), file_path
    except Exception as e:
        return f"Failed to process XML/GAN file: {str(e)}", None

def handle_audio_file(file_path):
    """Converts audio to text using the SpeechRecognition library."""
    if file_path.lower().endswith('.m4a'):
        sound = AudioSegment.from_file(file_path, format='m4a')
        wav_path = file_path.replace('.m4a', '.wav')
        sound.export(wav_path, format='wav')
        file_path = wav_path  # Update file path to the converted wav file

    recognizer = sr.Recognizer()
    with sr.AudioFile(file_path) as source:
        audio_data = recognizer.record(source)
        try:
            return recognizer.recognize_google(audio_data, language='it-IT'), file_path
        except sr.UnknownValueError:
            return "Speech not understood", file_path
        except sr.RequestError as e:
            return f"Speech recognition request failed; {e}", file_path

def handle_generic_video_file(file_path):
    """Handles video files for speech recognition from extracted audio."""
    audio_path = extract_audio_from_video(file_path)
    text = transcribe_audio(audio_path)
    os.remove(audio_path)  # Cleanup temporary audio file
    return text, file_path

def extract_audio_from_video(video_path):
    """Extracts the audio track from a video file and temporarily saves it as a WAV file."""
    video = VideoFileClip(video_path)
    audio_path = os.path.join(temp_directory, "temp_audio.wav")
    video.audio.write_audiofile(audio_path)
    return audio_path

def transcribe_audio(audio_path):
    """Converts audio to text using Google Speech Recognition, setting Italian as the default language."""
    recognizer = sr.Recognizer()
    with sr.AudioFile(audio_path) as source:
        audio_data = recognizer.record(source)  # Record audio from the entire file
        try:
            # Use 'en-IT' to set Italian as the recognition language
            return recognizer.recognize_google(audio_data, language='it-IT')
        except sr.UnknownValueError:
            return "Speech not understood"
        except sr.RequestError as e:
            return f"Could not request results; {e}"

def write_to_output(content, output_dir, file_index, original_path):
    """Writes the processed content to an output file, noting the path to the original file."""
    output_file_path = os.path.join(output_dir, f'modello_{file_index}.txt')
    with open(output_file_path, 'a', encoding='utf-8') as file:
        file.write(f"Original file path: {original_path}\n{content}\n")
    return file_index + 1  # Incrementa l'indice per il prossimo file

def handle_zip_file(zip_path):
    """Manages files within a ZIP archive, extracting and analyzing the contained files."""
    with zipfile.ZipFile(zip_path, 'r') as z:
        z.extractall(temp_directory)
        extracted_files = z.namelist()
        for file_name in extracted_files:
            internal_path = os.path.join(temp_directory, file_name)
            if os.path.isfile(internal_path):
                content, _ = handle_file(internal_path)
                if content and not content.startswith("Unsupported"):
                    return f"{content} (from {file_name} in {zip_path})", zip_path
                os.remove(internal_path)  # Clean up extracted files
    return "No supported files found or failed to process", None

def explore_directory(directory, output_dir, ignore_dirs, process_subfolders):
    """It explores directories and processes files according to specifications, excluding ignored folders."""
    file_index = 1  # Initialize the index of the output file
    for root, dirs, files in os.walk(directory):
        dirs[:] = [d for d in dirs if os.path.join(root, d) not in ignore_dirs]
        if not process_subfolders:
            dirs[:] = []
        for file in files:
            file_path = os.path.join(root, file)
            if any(os.path.abspath(os.path.join(root, d)) in ignore_dirs for d in dirs):
                continue
            content, original_path = handle_file(file_path)
            if content and not content.startswith("Unsupported"):
                file_index = write_to_output(content, output_dir, file_index, original_path)
            else:
                print(content)

def handle_file(file_path):
    """Determines the file type and applies the appropriate management function."""
    extension = os.path.splitext(file_path)[1].lower()
    handler = {
        '.txt': handle_text_file,
        '.htm': handle_epub_file,
        '.html': handle_epub_file,
        '.pdf': handle_pdf_file,
        '.docx': handle_word_file,
        '.doc': handle_word_file,
        '.pptx': handle_ppt_file,
        '.ppt': handle_ppt_file,
        '.xls': handle_excel_file,
        '.xlsx': handle_excel_file,
        '.xml': handle_xml_gan_file,
        '.gan': handle_xml_gan_file,
        '.xsd': handle_xml_gan_file,
        '.wav': handle_audio_file,
        '.mp3': handle_audio_file,
        '.m4a': handle_audio_file,
        '.mp4': handle_generic_video_file,
        '.avi': handle_generic_video_file,
        '.mov': handle_generic_video_file,
        '.mkv': handle_generic_video_file,
        '.mpeg': handle_generic_video_file,
        '.mpg': handle_generic_video_file,
        '.3gp': handle_generic_video_file,
        '.csv': handle_csv_file,
        '.epub': handle_epub_file,
        '.zip': handle_zip_file
    }.get(extension)
    if handler:
        return handler(file_path)
    return "Unsupported file format for {}".format(file_path), None

if __name__ == "__main__":
    directories = [r'C:\\someFolder']
    output_path = r'E:\\LoRA'
    ignore_dirs = [r'L:\\AFolderTobeIgnored\\', r'L:\\someOther']
    process_subfolders = True
    for directory in directories:
        explore_directory(directory, output_path, ignore_dirs, process_subfolders)

